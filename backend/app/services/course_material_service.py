from __future__ import annotations

import mimetypes
import re
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any

from bson import ObjectId
from fastapi import HTTPException

from app.chroma.material_store import ChromaMaterialStore, MaterialChunk
from app.core.config import settings
from app.repositories.course_material_repository import CourseMaterialRepository
from app.repositories.course_repository import CourseRepository
from app.repositories.user_repository import UserRepository
from app.schemas.course_material import CourseMaterialListResponse, CourseMaterialOut, CourseMaterialViewUrlResponse
from app.file_storage.r2 import R2Storage


MAX_CHUNK_WORDS = 450
CHUNK_OVERLAP_WORDS = 80


def _safe_filename(file_name: str) -> str:
    name = Path(file_name or "material").name.strip() or "material"
    return re.sub(r"[^A-Za-z0-9._ -]+", "_", name).strip(" .") or "material"


def _material_to_out(doc: dict[str, Any]) -> CourseMaterialOut:
    return CourseMaterialOut(
        id=str(doc["_id"]),
        course_id=str(doc["course_id"]),
        user_id=str(doc["user_id"]),
        file_name=str(doc["file_name"]),
        mime_type=str(doc["mime_type"]),
        file_size=int(doc["file_size"]),
        status=str(doc["status"]),
        created_at=doc["created_at"],
        updated_at=doc["updated_at"],
    )


def _chunk_words(*, material_id: str, page: int, text: str) -> list[MaterialChunk]:
    words = text.split()
    if not words:
        return []
    chunks: list[MaterialChunk] = []
    start = 0
    chunk_index = 0
    while start < len(words):
        end = min(start + MAX_CHUNK_WORDS, len(words))
        chunk_text = " ".join(words[start:end]).strip()
        if chunk_text:
            chunks.append(
                MaterialChunk(
                    chunk_id=f"{material_id}:p{page}:c{chunk_index}",
                    text=chunk_text,
                    page=page,
                    chunk_index=chunk_index,
                )
            )
            chunk_index += 1
        if end >= len(words):
            break
        start = max(end - CHUNK_OVERLAP_WORDS, start + 1)
    return chunks


def _extract_pdf_pages(file_bytes: bytes) -> list[tuple[int, str]]:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(file_bytes))
    pages: list[tuple[int, str]] = []
    for index, page in enumerate(reader.pages, start=1):
        pages.append((index, (page.extract_text() or "").strip()))
    return pages


def _extract_docx_pages(file_bytes: bytes) -> list[tuple[int, str]]:
    from docx import Document

    document = Document(BytesIO(file_bytes))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return [(1, "\n".join(parts).strip())]


def _extract_pptx_pages(file_bytes: bytes) -> list[tuple[int, str]]:
    from pptx import Presentation

    presentation = Presentation(BytesIO(file_bytes))
    pages: list[tuple[int, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                parts.append(text.strip())
        pages.append((index, "\n".join(parts).strip()))
    return pages


def _extract_text_pages(*, file_name: str, mime_type: str, file_bytes: bytes) -> list[tuple[int, str]]:
    suffix = Path(file_name).suffix.lower()
    if mime_type == "application/pdf" or suffix == ".pdf":
        return _extract_pdf_pages(file_bytes)
    if suffix == ".docx" or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _extract_docx_pages(file_bytes)
    if suffix == ".pptx" or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _extract_pptx_pages(file_bytes)
    if suffix in {".txt", ".md"} or mime_type.startswith("text/"):
        try:
            return [(1, file_bytes.decode("utf-8-sig").strip())]
        except UnicodeDecodeError:
            return [(1, file_bytes.decode("latin-1").strip())]
    raise ValueError("Unsupported file type. Upload PDF, TXT, MD, DOCX, or PPTX.")


class CourseMaterialService:
    def __init__(
        self,
        *,
        course_material_repo: CourseMaterialRepository,
        course_repo: CourseRepository,
        user_repo: UserRepository,
        r2_storage: R2Storage,
        chroma_material_store: ChromaMaterialStore,
    ):
        self.course_material_repo = course_material_repo
        self.course_repo = course_repo
        self.user_repo = user_repo
        self.r2_storage = r2_storage
        self.chroma_material_store = chroma_material_store

    def _authorize_professor_course(self, *, course_id: str, user_id: str, user_type: str | None) -> dict[str, Any]:
        if user_type != "professor":
            raise HTTPException(status_code=403, detail="Only professor accounts can access this resource")
        if not ObjectId.is_valid(course_id):
            raise HTTPException(status_code=422, detail="Invalid course_id")
        if not ObjectId.is_valid(user_id):
            raise HTTPException(status_code=401, detail="Invalid user token")
        course_doc = self.course_repo.get_by_id(course_id)
        if not course_doc:
            raise HTTPException(status_code=404, detail="Course not found")
        professor_institutions = set(self.user_repo.list_institution_ids(user_id=user_id, role="professor"))
        if str(course_doc.get("institution_id")) not in professor_institutions:
            raise HTTPException(status_code=403, detail="You do not have access to this course")
        return course_doc

    def list_materials(
        self,
        *,
        course_id: str,
        user_id: str,
        user_type: str | None,
    ) -> CourseMaterialListResponse:
        self._authorize_professor_course(course_id=course_id, user_id=user_id, user_type=user_type)
        rows = self.course_material_repo.list_by_course(course_id=ObjectId(course_id))
        return CourseMaterialListResponse(items=[_material_to_out(row) for row in rows])

    def create_material_upload(
        self,
        *,
        course_id: str,
        user_id: str,
        user_type: str | None,
        file_name: str,
        mime_type: str | None,
        file_bytes: bytes,
    ) -> CourseMaterialOut:
        self._authorize_professor_course(course_id=course_id, user_id=user_id, user_type=user_type)
        if not self.r2_storage.enabled:
            raise HTTPException(status_code=503, detail="Cloudflare R2 is not configured")
        if not self.chroma_material_store.enabled:
            raise HTTPException(status_code=503, detail="Chroma material collection is not configured")
        if not file_bytes:
            raise HTTPException(status_code=422, detail="Uploaded file is empty")

        material_id = ObjectId()
        safe_name = _safe_filename(file_name)
        content_type = (mime_type or mimetypes.guess_type(safe_name)[0] or "application/octet-stream").strip()
        storage_key = f"course-materials/{course_id}/{material_id}/{safe_name}"
        now = datetime.now(timezone.utc)

        try:
            self.r2_storage.upload_bytes(key=storage_key, body=file_bytes, content_type=content_type)
        except Exception as exc:  # noqa: BLE001
            raise HTTPException(status_code=502, detail=f"Failed to upload file to R2: {exc}") from exc

        doc = self.course_material_repo.create_material(
            material_id=material_id,
            course_id=ObjectId(course_id),
            user_id=ObjectId(user_id),
            file_name=safe_name,
            mime_type=content_type,
            file_size=len(file_bytes),
            storage_key=storage_key,
            status="processing",
            created_at=now,
            updated_at=now,
        )
        return _material_to_out(doc)

    def process_material_chunks(
        self,
        *,
        material_id: str,
        course_id: str,
        file_name: str,
        mime_type: str,
        file_bytes: bytes,
    ) -> None:
        material_oid = ObjectId(material_id)
        try:
            pages = _extract_text_pages(file_name=file_name, mime_type=mime_type, file_bytes=file_bytes)
            chunks: list[MaterialChunk] = []
            for page, text in pages:
                chunks.extend(_chunk_words(material_id=material_id, page=page, text=text))
            self.chroma_material_store.upsert_chunks(
                course_id=course_id,
                material_id=material_id,
                file_name=file_name,
                chunks=chunks,
            )
            self.course_material_repo.update_status(
                material_id=material_oid,
                status="ready",
                updated_at=datetime.now(timezone.utc),
            )
            print(
                f"[TA-BACKEND][materials] indexed material={material_id} chunks={len(chunks)}",
                flush=True,
            )
        except Exception as exc:  # noqa: BLE001
            self.course_material_repo.update_status(
                material_id=material_oid,
                status="failed",
                updated_at=datetime.now(timezone.utc),
            )
            print(f"[TA-BACKEND][materials] failed material={material_id}: {exc}", flush=True)

    def get_view_url(
        self,
        *,
        material_id: str,
        user_id: str,
        user_type: str | None,
    ) -> CourseMaterialViewUrlResponse:
        material = self.course_material_repo.get_by_id(material_id)
        if not material:
            raise HTTPException(status_code=404, detail="Material not found")
        self._authorize_professor_course(
            course_id=str(material["course_id"]),
            user_id=user_id,
            user_type=user_type,
        )
        if not self.r2_storage.enabled:
            raise HTTPException(status_code=503, detail="Cloudflare R2 is not configured")
        url = self.r2_storage.create_presigned_get_url(
            key=str(material["storage_key"]),
            expires_in=settings.r2_presigned_url_expires_seconds,
        )
        return CourseMaterialViewUrlResponse(url=url)
